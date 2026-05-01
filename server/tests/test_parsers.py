import json
import hashlib
import os
import re
import time
from io import BytesIO
from pathlib import Path
import zipfile

from docx import Document
from jsonschema import Draft202012Validator
import pytest
from pptx.enum.shapes import MSO_SHAPE
from pptx import Presentation
from pptx.util import Inches
from pypdf import PdfWriter
from pypdf.generic import DictionaryObject, NameObject, StreamObject

from server.ai.ocr import OcrAssetResult, OcrBox, VivoOcrClient
from server.ai.vision import VivoVisionClient, VisionAssetResult, VisionResult, VisualAsset
from server.parsers import DocxParser, ParserResult, PdfParser, PptxParser, SrtParser, VideoParser, parse_resource
from server.parsers.base import clean_text, text_quality_issue


ROOT = Path(__file__).resolve().parents[2]
NORMALIZED_DOCUMENT_SCHEMA = json.loads(
    (ROOT / "schemas/parse/normalized_document.schema.json").read_text(encoding="utf-8")
)
NORMALIZED_DOCUMENT_VALIDATOR = Draft202012Validator(NORMALIZED_DOCUMENT_SCHEMA)
PNG_BYTES = (
    b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01"
    b"\x08\x02\x00\x00\x00\x90wS\xde\x00\x00\x00\x0cIDATx\x9cc```\x00\x00"
    b"\x00\x04\x00\x01\xf6\x178U\x00\x00\x00\x00IEND\xaeB`\x82"
)


def assert_valid_normalized_document(payload: dict[str, object] | None) -> dict[str, object]:
    assert payload is not None
    NORMALIZED_DOCUMENT_VALIDATOR.validate(payload)
    return payload


def issue_codes(result) -> list[str]:
    return [issue.code for issue in result.issues]


def compact_text(text: str) -> str:
    return re.sub(r"[\W_]+", "", clean_text(text).lower(), flags=re.UNICODE)


class FakeVisionClient:
    def __init__(self, results: list[VisionResult] | dict[str, list[VisionResult]]) -> None:
        self.results = results
        self.calls: list[dict[str, object]] = []
        self.batch_calls: list[dict[str, object]] = []

    def analyze_image(self, image_bytes, *, mime_type, resource_type, location, hint=None):
        self.calls.append(
            {
                "image_bytes": image_bytes,
                "mime_type": mime_type,
                "resource_type": resource_type,
                "location": location,
                "hint": hint,
            }
        )
        return self._results_for("image-1")

    def analyze_images(self, assets, *, resource_type, document_context=None):
        self.batch_calls.append(
            {
                "assets": list(assets),
                "resource_type": resource_type,
                "document_context": document_context,
            }
        )
        results: list[VisionAssetResult] = []
        for asset in assets:
            self.calls.append(
                {
                    "image_bytes": asset.image_bytes,
                    "mime_type": asset.mime_type,
                    "resource_type": resource_type,
                    "location": asset.location,
                    "hint": asset.hint,
                }
            )
            for result in self._results_for(asset.asset_id):
                results.append(
                    VisionAssetResult(
                        asset_id=asset.asset_id,
                        segment_type=result.segment_type,
                        text=result.text,
                    )
                )
        return results

    def _results_for(self, asset_id: str) -> list[VisionResult]:
        if isinstance(self.results, dict):
            return self.results.get(asset_id, [])
        return self.results


class FakeOcrClient:
    def __init__(self, results: list[OcrAssetResult] | dict[str, list[OcrAssetResult]]) -> None:
        self.results = results
        self.calls: list[dict[str, object]] = []

    def recognize_images(self, assets, *, resource_type):
        self.calls.append({"assets": list(assets), "resource_type": resource_type})
        output: list[OcrAssetResult] = []
        for asset in assets:
            if isinstance(self.results, dict):
                output.extend(self.results.get(asset.asset_id, []))
            else:
                output.extend(
                    OcrAssetResult(asset_id=asset.asset_id, text=result.text, boxes=result.boxes)
                    for result in self.results
                )
        return output


class RaisingOcrClient:
    def recognize_images(self, assets, *, resource_type):
        raise RuntimeError("ocr unavailable")


class RaisingVisionClient:
    def analyze_image(self, image_bytes, *, mime_type, resource_type, location, hint=None):
        raise RuntimeError("vision unavailable")

    def analyze_images(self, assets, *, resource_type, document_context=None):
        raise RuntimeError("vision unavailable")


def create_text_pdf(path: Path, text: str) -> None:
    writer = PdfWriter()
    page = writer.add_blank_page(width=612, height=792)
    font = DictionaryObject(
        {
            NameObject("/Type"): NameObject("/Font"),
            NameObject("/Subtype"): NameObject("/Type1"),
            NameObject("/BaseFont"): NameObject("/Helvetica"),
        }
    )
    font_ref = writer._add_object(font)
    page[NameObject("/Resources")] = DictionaryObject(
        {NameObject("/Font"): DictionaryObject({NameObject("/F1"): font_ref})}
    )

    escaped = text.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")
    content = StreamObject()
    content._data = f"BT /F1 18 Tf 72 720 Td ({escaped}) Tj ET".encode("latin-1")
    page[NameObject("/Contents")] = writer._add_object(content)

    with path.open("wb") as file:
        writer.write(file)


def create_blank_pdf(path: Path) -> None:
    writer = PdfWriter()
    writer.add_blank_page(width=612, height=792)
    with path.open("wb") as file:
        writer.write(file)


def create_vector_venn_pdf(path: Path) -> None:
    import fitz

    document = fitz.open()
    page = document.new_page(width=612, height=792)
    page.insert_text((72, 72), "Venn diagram\nA\na", fontsize=14)
    page.draw_rect(fitz.Rect(90, 130, 240, 260), color=(0, 0, 0), width=1)
    page.draw_oval(fitz.Rect(120, 150, 210, 240), color=(0, 0, 0), width=1)
    document.save(path)
    document.close()


def create_docx_bytes(text: str) -> bytes:
    buffer = BytesIO()
    document = Document()
    document.add_paragraph(text)
    document.save(buffer)
    return buffer.getvalue()


def add_embedded_docx_relationship(pptx_path: Path, docx_bytes: bytes, *, slide_no: int = 1) -> None:
    rels_name = f"ppt/slides/_rels/slide{slide_no}.xml.rels"
    embedded_name = "ppt/embeddings/Embedded1.docx"
    content_types_name = "[Content_Types].xml"
    relationship = (
        '<Relationship Id="rId999" '
        'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/package" '
        'Target="../embeddings/Embedded1.docx"/>'
    )
    content_type = (
        '<Override PartName="/ppt/embeddings/Embedded1.docx" '
        'ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document"/>'
    )
    tmp_path = pptx_path.with_suffix(".embedded.pptx")
    with zipfile.ZipFile(pptx_path, "r") as source, zipfile.ZipFile(tmp_path, "w") as target:
        for item in source.infolist():
            data = source.read(item.filename)
            if item.filename == rels_name:
                text = data.decode("utf-8")
                data = text.replace("</Relationships>", f"{relationship}</Relationships>").encode("utf-8")
            elif item.filename == content_types_name:
                text = data.decode("utf-8")
                data = text.replace("</Types>", f"{content_type}</Types>").encode("utf-8")
            target.writestr(item, data)
        target.writestr(embedded_name, docx_bytes)
    tmp_path.replace(pptx_path)


def test_text_quality_removes_abnormal_script_noise():
    assert text_quality_issue("print() படட்") == "garbled"
    assert clean_text("print() படட்") == "print()"


def test_pdf_parser_extracts_page_text_and_validates_schema(tmp_path: Path):
    pdf_path = tmp_path / "text.pdf"
    create_text_pdf(pdf_path, "PDF Parser Text")

    result = PdfParser().parse(pdf_path)

    assert result.status == "succeeded"
    document = assert_valid_normalized_document(result.normalized_document)
    assert document["resourceType"] == "pdf"
    assert document["segments"] == [
        {
            "segmentKey": "pdf-p1",
            "segmentType": "pdf_page_text",
            "orderNo": 1,
            "textContent": "PDF Parser Text",
            "pageNo": 1,
        }
    ]


def test_pdf_parser_reports_empty_text_page(tmp_path: Path):
    pdf_path = tmp_path / "blank.pdf"
    create_blank_pdf(pdf_path)

    result = PdfParser().parse(pdf_path)

    assert result.status == "failed"
    assert result.normalized_document is None
    assert issue_codes(result) == ["pdf.page_text_empty"]


def test_pdf_parser_uses_ocr_fallback_for_garbled_page(monkeypatch, tmp_path: Path):
    class FakePage:
        def extract_text(self):
            return "bad\uffff\x00\x01\x19text"

    class FakeReader:
        pages = [FakePage()]

    monkeypatch.setattr("server.parsers.pdf.PdfReader", lambda _: FakeReader())
    pdf_path = tmp_path / "garbled.pdf"
    pdf_path.write_bytes(b"%PDF fake")
    vision_client = FakeVisionClient([VisionResult(segment_type="ocr_text", text="OCR clean text")])

    result = PdfParser(vision_client=vision_client).parse(pdf_path)

    assert result.status == "succeeded"
    document = assert_valid_normalized_document(result.normalized_document)
    assert vision_client.calls[0]["location"] == {"pageNo": 1}
    assert document["segments"] == [
        {
            "segmentKey": "pdf-p1-ocr1",
            "segmentType": "ocr_text",
            "orderNo": 1,
            "textContent": "OCR clean text",
            "pageNo": 1,
        }
    ]
    assert "\uffff" not in document["segments"][0]["textContent"]
    assert "\x00" not in document["segments"][0]["textContent"]


def test_pdf_parser_does_not_emit_garbled_text_when_ocr_is_not_configured(monkeypatch, tmp_path: Path):
    class FakePage:
        def extract_text(self):
            return "bad\uffff\x00\x01\x19text"

    class FakeReader:
        pages = [FakePage()]

    monkeypatch.setattr("server.parsers.pdf.PdfReader", lambda _: FakeReader())
    pdf_path = tmp_path / "garbled.pdf"
    pdf_path.write_bytes(b"%PDF fake")

    result = PdfParser().parse(pdf_path)

    assert result.status == "failed"
    assert result.normalized_document is None
    assert issue_codes(result) == ["pdf.page_text_garbled"]


def test_pdf_parser_appends_ocr_for_mixed_text_and_visual_page(monkeypatch, tmp_path: Path):
    pdf_path = tmp_path / "mixed.pdf"
    create_text_pdf(pdf_path, "Title")
    monkeypatch.setattr("server.parsers.pdf._pymupdf_page_needs_visual_ocr", lambda page, text: True)
    monkeypatch.setattr("server.parsers.pdf._pymupdf_page_png", lambda page: b"page-png")
    vision_client = FakeVisionClient([VisionResult(segment_type="ocr_text", text="Screenshot body text")])

    result = PdfParser(vision_client=vision_client).parse(pdf_path)

    assert result.status == "succeeded"
    document = assert_valid_normalized_document(result.normalized_document)
    assert [segment["segmentType"] for segment in document["segments"]] == ["pdf_page_text", "ocr_text"]
    assert document["segments"][0]["pageNo"] == 1
    assert document["segments"][1] == {
        "segmentKey": "pdf-p1-ocr1",
        "segmentType": "ocr_text",
        "orderNo": 2,
        "textContent": "Screenshot body text",
        "pageNo": 1,
    }


def test_pdf_parser_skips_duplicate_mixed_page_ocr(monkeypatch, tmp_path: Path):
    pdf_path = tmp_path / "mixed.pdf"
    create_text_pdf(pdf_path, "Duplicate visible text")
    monkeypatch.setattr("server.parsers.pdf._pymupdf_page_needs_visual_ocr", lambda page, text: True)
    monkeypatch.setattr("server.parsers.pdf._pymupdf_page_png", lambda page: b"page-png")
    vision_client = FakeVisionClient([VisionResult(segment_type="ocr_text", text="Duplicate visible text")])

    result = PdfParser(vision_client=vision_client).parse(pdf_path)

    assert result.status == "succeeded"
    document = assert_valid_normalized_document(result.normalized_document)
    assert [segment["segmentType"] for segment in document["segments"]] == ["pdf_page_text"]


def test_pdf_parser_skips_short_duplicate_mixed_page_ocr(monkeypatch, tmp_path: Path):
    pdf_path = tmp_path / "mixed.pdf"
    create_text_pdf(pdf_path, "Title")
    monkeypatch.setattr("server.parsers.pdf._pymupdf_page_needs_visual_ocr", lambda page, text: True)
    monkeypatch.setattr("server.parsers.pdf._pymupdf_page_png", lambda page: b"page-png")
    vision_client = FakeVisionClient([VisionResult(segment_type="ocr_text", text="Title")])

    result = PdfParser(vision_client=vision_client).parse(pdf_path)

    assert result.status == "succeeded"
    document = assert_valid_normalized_document(result.normalized_document)
    assert [segment["segmentType"] for segment in document["segments"]] == ["pdf_page_text"]


def test_pdf_parser_reports_empty_visual_result_for_mixed_page(monkeypatch, tmp_path: Path):
    pdf_path = tmp_path / "mixed.pdf"
    create_text_pdf(pdf_path, "Title")
    monkeypatch.setattr("server.parsers.pdf._pymupdf_page_needs_visual_ocr", lambda page, text: True)
    monkeypatch.setattr("server.parsers.pdf._pymupdf_page_png", lambda page: b"page-png")

    result = PdfParser(vision_client=FakeVisionClient([])).parse(pdf_path)

    assert result.status == "succeeded"
    assert issue_codes(result) == ["pdf.visual_empty"]


def test_pdf_parser_keeps_text_when_mixed_page_ocr_fails(monkeypatch, tmp_path: Path):
    pdf_path = tmp_path / "mixed.pdf"
    create_text_pdf(pdf_path, "Title")
    monkeypatch.setattr("server.parsers.pdf._pymupdf_page_needs_visual_ocr", lambda page, text: True)
    monkeypatch.setattr("server.parsers.pdf._pymupdf_page_png", lambda page: b"page-png")

    result = PdfParser(vision_client=RaisingVisionClient()).parse(pdf_path)

    assert result.status == "succeeded"
    document = assert_valid_normalized_document(result.normalized_document)
    assert [segment["segmentType"] for segment in document["segments"]] == ["pdf_page_text"]
    assert issue_codes(result) == ["pdf.vision_failed"]


def test_pdf_parser_uses_vivo_ocr_before_vision(monkeypatch, tmp_path: Path):
    pdf_path = tmp_path / "ocr-first.pdf"
    create_text_pdf(pdf_path, "Title")
    monkeypatch.setattr("server.parsers.pdf._pymupdf_page_needs_visual_ocr", lambda page, text: True)
    monkeypatch.setattr("server.parsers.pdf._pymupdf_page_png", lambda page: b"page-png")
    ocr_text = "Screenshot body text with enough useful details for the OCR quality gate."
    ocr_client = FakeOcrClient([OcrAssetResult(asset_id="unused", text=ocr_text, boxes=[OcrBox(text=ocr_text)])])
    vision_client = FakeVisionClient([VisionResult(segment_type="image_caption", text="should not be called")])

    result = PdfParser(ocr_client=ocr_client, vision_client=vision_client).parse(pdf_path)

    assert result.status == "succeeded"
    document = assert_valid_normalized_document(result.normalized_document)
    assert len(ocr_client.calls) == 1
    assert [asset.asset_id for asset in ocr_client.calls[0]["assets"]] == ["pdf-p1"]
    assert vision_client.calls == []
    assert document["segments"][1] == {
        "segmentKey": "pdf-p1-ocr1",
        "segmentType": "ocr_text",
        "orderNo": 2,
        "textContent": ocr_text,
        "pageNo": 1,
    }
    assert "boxes" not in document["segments"][1]


def test_pdf_parser_calls_vision_when_ocr_is_low_quality(monkeypatch, tmp_path: Path):
    pdf_path = tmp_path / "ocr-low.pdf"
    create_text_pdf(pdf_path, "Title")
    monkeypatch.setattr("server.parsers.pdf._pymupdf_page_needs_visual_ocr", lambda page, text: True)
    monkeypatch.setattr("server.parsers.pdf._pymupdf_page_png", lambda page: b"page-png")
    ocr_client = FakeOcrClient([OcrAssetResult(asset_id="unused", text="x", boxes=[])])
    vision_client = FakeVisionClient([VisionResult(segment_type="ocr_text", text="Recovered visual text from VLM")])

    result = PdfParser(ocr_client=ocr_client, vision_client=vision_client).parse(pdf_path)

    assert result.status == "succeeded"
    document = assert_valid_normalized_document(result.normalized_document)
    assert len(ocr_client.calls) == 1
    assert len(vision_client.calls) == 1
    assert [segment["textContent"] for segment in document["segments"]] == [
        "Title",
        "Recovered visual text from VLM",
    ]


def test_pdf_parser_rejects_low_quality_ocr_without_vision_for_empty_text_layer(tmp_path: Path):
    pdf_path = tmp_path / "blank-with-broken-ocr.pdf"
    create_blank_pdf(pdf_path)
    broken_ocr = "常用结论\n1.子集有2”个\n2.等价关系：ABeAB=A→AUB=B→CuACuB."
    ocr_client = FakeOcrClient([OcrAssetResult(asset_id="unused", text=broken_ocr, boxes=[])])

    result = PdfParser(ocr_client=ocr_client).parse(pdf_path)

    assert result.status == "failed"
    assert result.normalized_document is None
    assert issue_codes(result) == ["pdf.page_text_empty"]


def test_pdf_parser_falls_back_to_vision_when_ocr_fails(monkeypatch, tmp_path: Path):
    pdf_path = tmp_path / "ocr-failed.pdf"
    create_text_pdf(pdf_path, "Title")
    monkeypatch.setattr("server.parsers.pdf._pymupdf_page_needs_visual_ocr", lambda page, text: True)
    monkeypatch.setattr("server.parsers.pdf._pymupdf_page_png", lambda page: b"page-png")
    vision_client = FakeVisionClient([VisionResult(segment_type="ocr_text", text="Vision fallback text")])

    result = PdfParser(ocr_client=RaisingOcrClient(), vision_client=vision_client).parse(pdf_path)

    assert result.status == "succeeded"
    document = assert_valid_normalized_document(result.normalized_document)
    assert len(vision_client.calls) == 1
    assert [segment["textContent"] for segment in document["segments"]] == ["Title", "Vision fallback text"]
    assert issue_codes(result) == ["pdf.ocr_failed"]


def test_pdf_parser_suppresses_broken_code_ocr_and_uses_vision(monkeypatch, tmp_path: Path):
    pdf_path = tmp_path / "code-flow.pdf"
    create_text_pdf(pdf_path, "print function explanation has enough native text to remain as the main layer.")
    monkeypatch.setattr("server.parsers.pdf._pymupdf_page_needs_visual_ocr", lambda page, text: True)
    monkeypatch.setattr("server.parsers.pdf._pymupdf_page_png", lambda page: b"page-png")
    broken_ocr = "print('The quick brown fox'\n'jumps over!\nlazy\nthe"
    ocr_client = FakeOcrClient([OcrAssetResult(asset_id="unused", text=broken_ocr, boxes=[])])
    vision_client = FakeVisionClient([VisionResult(segment_type="image_caption", text="流程图说明 print 会用空格连接参数。")])

    result = PdfParser(ocr_client=ocr_client, vision_client=vision_client).parse(pdf_path)

    assert result.status == "succeeded"
    document = assert_valid_normalized_document(result.normalized_document)
    assert len(vision_client.calls) == 1
    assert [segment["segmentType"] for segment in document["segments"]] == ["pdf_page_text", "image_caption"]
    assert all(broken_ocr not in segment["textContent"] for segment in document["segments"])


def test_pdf_parser_skips_cover_visual_when_ocr_empty(monkeypatch, tmp_path: Path):
    pdf_path = tmp_path / "cover.pdf"
    create_text_pdf(pdf_path, "Course Title")
    monkeypatch.setattr("server.parsers.pdf._pymupdf_page_needs_visual_ocr", lambda page, text: True)
    monkeypatch.setattr("server.parsers.pdf._pymupdf_page_png", lambda page: b"page-png")
    ocr_client = FakeOcrClient([])
    vision_client = FakeVisionClient([VisionResult(segment_type="image_caption", text="cover caption")])

    result = PdfParser(ocr_client=ocr_client, vision_client=vision_client).parse(pdf_path)

    assert result.status == "succeeded"
    document = assert_valid_normalized_document(result.normalized_document)
    assert vision_client.calls == []
    assert document["segments"] == [
        {
            "segmentKey": "pdf-p1",
            "segmentType": "pdf_page_text",
            "orderNo": 1,
            "textContent": "Course Title",
            "pageNo": 1,
        }
    ]
    assert issue_codes(result) == []


def test_pdf_parser_adds_local_caption_for_vector_venn_without_vision(tmp_path: Path):
    pdf_path = tmp_path / "venn-vector.pdf"
    create_vector_venn_pdf(pdf_path)

    result = PdfParser().parse(pdf_path)

    assert result.status == "succeeded"
    document = assert_valid_normalized_document(result.normalized_document)
    assert issue_codes(result) == []
    assert [segment["segmentType"] for segment in document["segments"]] == ["pdf_page_text", "image_caption"]
    assert document["segments"][1]["pageNo"] == 1
    assert "文氏图" in document["segments"][1]["textContent"]
    assert "元素 a" in document["segments"][1]["textContent"]


def test_pptx_parser_extracts_slide_text_and_validates_schema(tmp_path: Path):
    pptx_path = tmp_path / "deck.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    textbox.text = "Slide Title\nSlide body"
    presentation.save(pptx_path)

    result = PptxParser().parse(pptx_path)

    assert result.status == "succeeded"
    document = assert_valid_normalized_document(result.normalized_document)
    assert document["resourceType"] == "pptx"
    assert document["segments"][0]["segmentKey"] == "pptx-s1"
    assert document["segments"][0]["segmentType"] == "ppt_slide_text"
    assert document["segments"][0]["slideNo"] == 1
    assert document["segments"][0]["textContent"] == "Slide Title\nSlide body"


def test_pptx_parser_adds_local_caption_for_vector_venn_without_vision(tmp_path: Path):
    pptx_path = tmp_path / "venn-vector.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    title = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(3), Inches(0.5))
    title.text = "文氏图"
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.3), Inches(2.2), Inches(1.4))
    slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.3), Inches(1.5), Inches(1.1), Inches(1.1))
    label_a = slide.shapes.add_textbox(Inches(1.55), Inches(1.75), Inches(0.4), Inches(0.3))
    label_a.text = "A"
    point_a = slide.shapes.add_textbox(Inches(2.3), Inches(2.1), Inches(0.4), Inches(0.3))
    point_a.text = "a"
    presentation.save(pptx_path)

    result = PptxParser().parse(pptx_path)

    assert result.status == "succeeded"
    document = assert_valid_normalized_document(result.normalized_document)
    assert issue_codes(result) == []
    assert [segment["segmentType"] for segment in document["segments"]] == ["ppt_slide_text", "image_caption"]
    assert document["segments"][1]["slideNo"] == 1
    assert "文氏图" in document["segments"][1]["textContent"]
    assert "元素 a" in document["segments"][1]["textContent"]


def test_pptx_parser_extracts_embedded_docx_text_and_skips_render_issue(monkeypatch, tmp_path: Path):
    pptx_path = tmp_path / "embedded-docx.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
    textbox.text = "Native slide text"
    presentation.save(pptx_path)
    add_embedded_docx_relationship(
        pptx_path,
        create_docx_bytes("2.A⊆B(子集)A⫋B(真子集)⇔A⊆B且A≠B,A=B(相等)⇔A⊆B且B⊆A."),
    )
    monkeypatch.setattr("server.parsers.pptx._slide_needs_render_fallback", lambda shapes: True)
    monkeypatch.setattr("server.parsers.pptx._render_slide_png", lambda file_path, slide_no: None)

    result = PptxParser(vision_client=FakeVisionClient([])).parse(pptx_path)

    assert result.status == "succeeded"
    document = assert_valid_normalized_document(result.normalized_document)
    assert document["segments"][0]["textContent"] == (
        "Native slide text\n2.A⊆B(子集)A⫋B(真子集)⇔A⊆B且A≠B,A=B(相等)⇔A⊆B且B⊆A."
    )
    assert "pptx.slide_render_unavailable" not in issue_codes(result)


def test_pptx_parser_reports_empty_slide(tmp_path: Path):
    pptx_path = tmp_path / "blank.pptx"
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[6])
    presentation.save(pptx_path)

    result = PptxParser().parse(pptx_path)

    assert result.status == "failed"
    assert result.normalized_document is None
    assert issue_codes(result) == ["pptx.slide_text_empty"]


def test_pptx_parser_merges_table_overlay_text(tmp_path: Path):
    pptx_path = tmp_path / "table-overlay.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    table_shape = slide.shapes.add_table(2, 3, Inches(1), Inches(1), Inches(6), Inches(2))
    table = table_shape.table
    table.cell(0, 0).text = "运算"
    table.cell(0, 1).text = "自然语言"
    table.cell(0, 2).text = "符号语言"
    table.cell(1, 0).text = "交集"
    table.cell(1, 1).text = "由所有属于A且属于B的元素组成"
    table.cell(1, 2).text = "A∩B="
    overlay = slide.shapes.add_textbox(Inches(4.5), Inches(2), Inches(1.5), Inches(0.4))
    overlay.text = "{x|x∈A,且x∈B}"
    presentation.save(pptx_path)

    result = PptxParser().parse(pptx_path)

    assert result.status == "succeeded"
    document = assert_valid_normalized_document(result.normalized_document)
    assert "A∩B= {x|x∈A,且x∈B}" in document["segments"][0]["textContent"]


def test_pptx_parser_extracts_image_formula_with_slide_location(tmp_path: Path):
    pptx_path = tmp_path / "formula.pptx"
    image_path = tmp_path / "formula.png"
    image_path.write_bytes(PNG_BYTES)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture(str(image_path), Inches(1), Inches(1), width=Inches(1))
    presentation.save(pptx_path)
    vision_client = FakeVisionClient([VisionResult(segment_type="formula", text="E = mc^2")])

    result = PptxParser(vision_client=vision_client).parse(pptx_path)

    assert result.status == "succeeded"
    document = assert_valid_normalized_document(result.normalized_document)
    assert vision_client.calls[0]["location"] == {"slideNo": 1, "imageNo": 1}
    assert document["segments"] == [
        {
            "segmentKey": "pptx-s1-i1-formula1",
            "segmentType": "formula",
            "orderNo": 1,
            "textContent": "E = mc^2",
            "slideNo": 1,
        }
    ]


def test_pptx_parser_batches_images_for_same_file(monkeypatch, tmp_path: Path):
    monkeypatch.setenv("KNOWLINK_VIVO_VISION_BATCH_SIZE", "5")
    pptx_path = tmp_path / "images.pptx"
    image_path = tmp_path / "diagram.png"
    image_path.write_bytes(PNG_BYTES)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture(str(image_path), Inches(1), Inches(1), width=Inches(1))
    slide.shapes.add_picture(str(image_path), Inches(3), Inches(1), width=Inches(1))
    presentation.save(pptx_path)
    vision_client = FakeVisionClient(
        {
            "pptx-s1-i1": [VisionResult(segment_type="image_caption", text="集合图 A")],
            "pptx-s1-i2": [VisionResult(segment_type="image_caption", text="集合图 B")],
        }
    )

    result = PptxParser(vision_client=vision_client).parse(pptx_path)

    assert result.status == "succeeded"
    document = assert_valid_normalized_document(result.normalized_document)
    assert len(vision_client.batch_calls) == 1
    assert [asset.asset_id for asset in vision_client.batch_calls[0]["assets"]] == ["pptx-s1-i1", "pptx-s1-i2"]
    assert [segment["textContent"] for segment in document["segments"]] == ["集合图 A", "集合图 B"]


def test_pptx_parser_keeps_slide_text_and_appends_image_ocr(tmp_path: Path):
    pptx_path = tmp_path / "mixed.pptx"
    image_path = tmp_path / "screenshot.png"
    image_path.write_bytes(PNG_BYTES)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    textbox.text = "Native slide text"
    slide.shapes.add_picture(str(image_path), Inches(1), Inches(2), width=Inches(1))
    presentation.save(pptx_path)
    vision_client = FakeVisionClient([VisionResult(segment_type="ocr_text", text="Screenshot text")])

    result = PptxParser(vision_client=vision_client).parse(pptx_path)

    assert result.status == "succeeded"
    document = assert_valid_normalized_document(result.normalized_document)
    assert document["segments"] == [
        {
            "segmentKey": "pptx-s1",
            "segmentType": "ppt_slide_text",
            "orderNo": 1,
            "textContent": "Native slide text",
            "slideNo": 1,
        },
        {
            "segmentKey": "pptx-s1-i1-ocr1",
            "segmentType": "ocr_text",
            "orderNo": 2,
            "textContent": "Screenshot text",
            "slideNo": 1,
        },
    ]


def test_pptx_parser_uses_vivo_ocr_before_vision(tmp_path: Path):
    pptx_path = tmp_path / "ocr-first.pptx"
    image_path = tmp_path / "screenshot.png"
    image_path.write_bytes(PNG_BYTES)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture(str(image_path), Inches(1), Inches(1), width=Inches(1))
    presentation.save(pptx_path)
    ocr_text = "Screenshot text with enough useful details for OCR quality gate."
    ocr_client = FakeOcrClient([OcrAssetResult(asset_id="unused", text=ocr_text, boxes=[OcrBox(text=ocr_text)])])
    vision_client = FakeVisionClient([VisionResult(segment_type="image_caption", text="should not be called")])

    result = PptxParser(ocr_client=ocr_client, vision_client=vision_client).parse(pptx_path)

    assert result.status == "succeeded"
    document = assert_valid_normalized_document(result.normalized_document)
    assert len(ocr_client.calls) == 1
    assert [asset.asset_id for asset in ocr_client.calls[0]["assets"]] == ["pptx-s1-i1"]
    assert vision_client.calls == []
    assert document["segments"] == [
        {
            "segmentKey": "pptx-s1-i1-ocr1",
            "segmentType": "ocr_text",
            "orderNo": 1,
            "textContent": ocr_text,
            "slideNo": 1,
        }
    ]
    assert "boxes" not in document["segments"][0]


def test_pptx_parser_calls_vision_when_ocr_is_low_quality(tmp_path: Path):
    pptx_path = tmp_path / "ocr-low.pptx"
    image_path = tmp_path / "diagram.png"
    image_path.write_bytes(PNG_BYTES)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture(str(image_path), Inches(1), Inches(1), width=Inches(1))
    presentation.save(pptx_path)
    ocr_client = FakeOcrClient([OcrAssetResult(asset_id="unused", text="A", boxes=[])])
    vision_client = FakeVisionClient([VisionResult(segment_type="image_caption", text="Venn 图展示集合关系。")])

    result = PptxParser(ocr_client=ocr_client, vision_client=vision_client).parse(pptx_path)

    assert result.status == "succeeded"
    document = assert_valid_normalized_document(result.normalized_document)
    assert len(ocr_client.calls) == 1
    assert len(vision_client.calls) == 1
    assert document["segments"][0]["segmentType"] == "image_caption"
    assert document["segments"][0]["textContent"] == "Venn 图展示集合关系。"


def test_pptx_parser_forces_vision_for_venn_label_ocr(tmp_path: Path):
    pptx_path = tmp_path / "venn-label.pptx"
    image_path = tmp_path / "venn.png"
    image_path.write_bytes(PNG_BYTES)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture(str(image_path), Inches(1), Inches(1), width=Inches(1))
    presentation.save(pptx_path)
    ocr_client = FakeOcrClient([OcrAssetResult(asset_id="unused", text="B\nA\nANB", boxes=[])])
    vision_client = FakeVisionClient([VisionResult(segment_type="image_caption", text="Venn 图阴影区域表示 A∩B。")])

    result = PptxParser(ocr_client=ocr_client, vision_client=vision_client).parse(pptx_path)

    assert result.status == "succeeded"
    document = assert_valid_normalized_document(result.normalized_document)
    assert len(vision_client.calls) == 1
    assert document["segments"] == [
        {
            "segmentKey": "pptx-s1-i1-image1",
            "segmentType": "image_caption",
            "orderNo": 1,
            "textContent": "Venn 图阴影区域表示 A∩B。",
            "slideNo": 1,
        }
    ]


def test_pptx_parser_forces_vision_for_broken_math_ocr(tmp_path: Path):
    pptx_path = tmp_path / "broken-math.pptx"
    image_path = tmp_path / "math.png"
    image_path.write_bytes(PNG_BYTES)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture(str(image_path), Inches(1), Inches(1), width=Inches(1))
    presentation.save(pptx_path)
    broken_ocr = "常用结论\n1.子集有2”个\n2.等价关系：ABeAB=A→AUB=B→CuACuB."
    ocr_client = FakeOcrClient([OcrAssetResult(asset_id="unused", text=broken_ocr, boxes=[])])
    vision_client = FakeVisionClient(
        [VisionResult(segment_type="formula", text="A⊆B⇔A∩B=A⇔A∪B=B⇔∁_U A⊇∁_U B")]
    )

    result = PptxParser(ocr_client=ocr_client, vision_client=vision_client).parse(pptx_path)

    assert result.status == "succeeded"
    document = assert_valid_normalized_document(result.normalized_document)
    assert len(vision_client.calls) == 1
    assert document["segments"][0]["segmentType"] == "formula"
    assert "AUB" not in document["segments"][0]["textContent"]
    assert "2”" not in document["segments"][0]["textContent"]


def test_pptx_parser_uses_vision_when_ocr_empty_even_with_slide_text(tmp_path: Path):
    pptx_path = tmp_path / "empty-ocr-with-text.pptx"
    image_path = tmp_path / "complement.png"
    image_path.write_bytes(PNG_BYTES)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    textbox.text = "3.集合的基本运算\n补集"
    slide.shapes.add_picture(str(image_path), Inches(1), Inches(2), width=Inches(1))
    presentation.save(pptx_path)
    ocr_client = FakeOcrClient([])
    vision_client = FakeVisionClient([VisionResult(segment_type="image_caption", text="灰色区域表示 ∁_U A。")])

    result = PptxParser(ocr_client=ocr_client, vision_client=vision_client).parse(pptx_path)

    assert result.status == "succeeded"
    document = assert_valid_normalized_document(result.normalized_document)
    assert len(vision_client.calls) == 1
    assert [segment["segmentType"] for segment in document["segments"]] == ["ppt_slide_text", "image_caption"]
    assert issue_codes(result) == []


def test_pptx_parser_skips_duplicate_image_ocr(tmp_path: Path):
    pptx_path = tmp_path / "duplicate.pptx"
    image_path = tmp_path / "screenshot.png"
    image_path.write_bytes(PNG_BYTES)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    textbox.text = "Duplicate slide text"
    slide.shapes.add_picture(str(image_path), Inches(1), Inches(2), width=Inches(1))
    presentation.save(pptx_path)
    vision_client = FakeVisionClient([VisionResult(segment_type="ocr_text", text="Duplicate slide text")])

    result = PptxParser(vision_client=vision_client).parse(pptx_path)

    assert result.status == "succeeded"
    document = assert_valid_normalized_document(result.normalized_document)
    assert [segment["segmentType"] for segment in document["segments"]] == ["ppt_slide_text"]


def test_pptx_parser_skips_short_duplicate_image_ocr(tmp_path: Path):
    pptx_path = tmp_path / "duplicate-short.pptx"
    image_path = tmp_path / "screenshot.png"
    image_path.write_bytes(PNG_BYTES)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    textbox.text = "语法知识"
    slide.shapes.add_picture(str(image_path), Inches(1), Inches(2), width=Inches(1))
    presentation.save(pptx_path)
    vision_client = FakeVisionClient([VisionResult(segment_type="ocr_text", text="语法知识")])

    result = PptxParser(vision_client=vision_client).parse(pptx_path)

    assert result.status == "succeeded"
    document = assert_valid_normalized_document(result.normalized_document)
    assert [segment["segmentType"] for segment in document["segments"]] == ["ppt_slide_text"]


def test_pptx_parser_reports_empty_visual_result(tmp_path: Path):
    pptx_path = tmp_path / "empty-visual.pptx"
    image_path = tmp_path / "screenshot.png"
    image_path.write_bytes(PNG_BYTES)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture(str(image_path), Inches(1), Inches(1), width=Inches(1))
    presentation.save(pptx_path)

    result = PptxParser(vision_client=FakeVisionClient([])).parse(pptx_path)

    assert result.status == "failed"
    assert issue_codes(result) == ["pptx.visual_empty"]


def test_pptx_parser_uses_slide_render_for_render_only_visual(monkeypatch, tmp_path: Path):
    pptx_path = tmp_path / "render-only.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    textbox.text = "Native slide text"
    presentation.save(pptx_path)
    monkeypatch.setattr("server.parsers.pptx._slide_needs_render_fallback", lambda shapes: True)
    monkeypatch.setattr("server.parsers.pptx._render_slide_png", lambda file_path, slide_no: b"slide-render")
    vision_client = FakeVisionClient(
        {"pptx-s1-render": [VisionResult(segment_type="ocr_text", text="Rendered OLE text")]}
    )

    result = PptxParser(vision_client=vision_client).parse(pptx_path)

    assert result.status == "succeeded"
    document = assert_valid_normalized_document(result.normalized_document)
    assert vision_client.calls[0]["image_bytes"] == b"slide-render"
    assert [segment["textContent"] for segment in document["segments"]] == [
        "Native slide text",
        "Rendered OLE text",
    ]


def test_pptx_parser_reports_unavailable_slide_render(monkeypatch, tmp_path: Path):
    pptx_path = tmp_path / "render-unavailable.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    textbox.text = "Native slide text"
    presentation.save(pptx_path)
    monkeypatch.setattr("server.parsers.pptx._slide_needs_render_fallback", lambda shapes: True)
    monkeypatch.setattr("server.parsers.pptx._render_slide_png", lambda file_path, slide_no: None)

    result = PptxParser(vision_client=FakeVisionClient([])).parse(pptx_path)

    assert result.status == "succeeded"
    assert issue_codes(result) == ["pptx.slide_render_unavailable"]


def test_pptx_parser_keeps_slide_text_when_image_ocr_fails(tmp_path: Path):
    pptx_path = tmp_path / "ocr-failed.pptx"
    image_path = tmp_path / "screenshot.png"
    image_path.write_bytes(PNG_BYTES)
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
    textbox.text = "Native slide text"
    slide.shapes.add_picture(str(image_path), Inches(1), Inches(2), width=Inches(1))
    presentation.save(pptx_path)

    result = PptxParser(vision_client=RaisingVisionClient()).parse(pptx_path)

    assert result.status == "succeeded"
    document = assert_valid_normalized_document(result.normalized_document)
    assert [segment["segmentType"] for segment in document["segments"]] == ["ppt_slide_text"]
    assert issue_codes(result) == ["pptx.vision_failed"]


def test_docx_parser_extracts_heading_section_path_and_validates_schema(tmp_path: Path):
    docx_path = tmp_path / "document.docx"
    document = Document()
    document.add_heading("Chapter 1", level=1)
    document.add_paragraph("First paragraph")
    document.add_heading("Section 1.1", level=2)
    document.add_paragraph("Second paragraph")
    document.save(docx_path)

    result = DocxParser().parse(docx_path)

    assert result.status == "succeeded"
    normalized = assert_valid_normalized_document(result.normalized_document)
    segments = normalized["segments"]
    assert segments[0]["segmentType"] == "docx_block_text"
    assert segments[0]["sectionPath"] == ["Chapter 1"]
    assert segments[1]["textContent"] == "First paragraph"
    assert segments[1]["sectionPath"] == ["Chapter 1"]
    assert segments[3]["textContent"] == "Second paragraph"
    assert segments[3]["sectionPath"] == ["Chapter 1", "Section 1.1"]


def test_docx_parser_uses_chinese_heading_heuristic_and_extracts_table(tmp_path: Path):
    docx_path = tmp_path / "table.docx"
    document = Document()
    document.add_paragraph("一、函数极限")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "概念"
    table.cell(0, 1).text = "定义"
    table.cell(1, 0).text = "极限"
    table.cell(1, 1).text = "趋近时的稳定趋势"
    document.save(docx_path)

    result = DocxParser().parse(docx_path)

    assert result.status == "succeeded"
    document_payload = assert_valid_normalized_document(result.normalized_document)
    segments = document_payload["segments"]
    assert segments[0]["sectionPath"] == ["一、函数极限"]
    assert segments[1] == {
        "segmentKey": "docx-b2",
        "segmentType": "docx_block_text",
        "orderNo": 2,
        "textContent": "概念 | 定义\n极限 | 趋近时的稳定趋势",
        "sectionPath": ["一、函数极限"],
    }


def test_docx_parser_extracts_image_caption_with_section_path(tmp_path: Path):
    docx_path = tmp_path / "image.docx"
    document = Document()
    document.add_heading("图表章节", level=1)
    document.add_picture(BytesIO(PNG_BYTES), width=Inches(1))
    document.save(docx_path)
    vision_client = FakeVisionClient([VisionResult(segment_type="image_caption", text="函数变化趋势图")])

    result = DocxParser(vision_client=vision_client).parse(docx_path)

    assert result.status == "succeeded"
    document_payload = assert_valid_normalized_document(result.normalized_document)
    assert vision_client.calls[0]["location"]["sectionPath"] == ["图表章节"]
    assert document_payload["segments"][1] == {
        "segmentKey": "docx-b2-i1-image1",
        "segmentType": "image_caption",
        "orderNo": 2,
        "textContent": "函数变化趋势图",
        "sectionPath": ["图表章节"],
    }


def test_docx_parser_uses_vivo_ocr_before_vision(tmp_path: Path):
    docx_path = tmp_path / "image-ocr.docx"
    document = Document()
    document.add_heading("图表章节", level=1)
    document.add_picture(BytesIO(PNG_BYTES), width=Inches(1))
    document.save(docx_path)
    ocr_text = "截图文字包含足够多的有效内容，会直接作为 OCR 文本进入解析结果，并保留题干、条件和关键结论。"
    ocr_client = FakeOcrClient([OcrAssetResult(asset_id="unused", text=ocr_text, boxes=[OcrBox(text=ocr_text)])])
    vision_client = FakeVisionClient([VisionResult(segment_type="image_caption", text="should not be called")])

    result = DocxParser(ocr_client=ocr_client, vision_client=vision_client).parse(docx_path)

    assert result.status == "succeeded"
    document_payload = assert_valid_normalized_document(result.normalized_document)
    assert len(ocr_client.calls) == 1
    assert [asset.asset_id for asset in ocr_client.calls[0]["assets"]] == ["docx-b2-i1"]
    assert vision_client.calls == []
    assert document_payload["segments"][1] == {
        "segmentKey": "docx-b2-i1-ocr1",
        "segmentType": "ocr_text",
        "orderNo": 2,
        "textContent": ocr_text,
        "sectionPath": ["图表章节"],
    }


def test_docx_parser_calls_vision_when_ocr_is_low_quality(tmp_path: Path):
    docx_path = tmp_path / "image-low-ocr.docx"
    document = Document()
    document.add_heading("图表章节", level=1)
    document.add_picture(BytesIO(PNG_BYTES), width=Inches(1))
    document.save(docx_path)
    ocr_client = FakeOcrClient([OcrAssetResult(asset_id="unused", text="B\nA\nANB", boxes=[])])
    vision_client = FakeVisionClient([VisionResult(segment_type="image_caption", text="文氏图阴影区域表示 A∩B。")])

    result = DocxParser(ocr_client=ocr_client, vision_client=vision_client).parse(docx_path)

    assert result.status == "succeeded"
    document_payload = assert_valid_normalized_document(result.normalized_document)
    assert len(ocr_client.calls) == 1
    assert len(vision_client.calls) == 1
    assert document_payload["segments"][1]["segmentType"] == "image_caption"
    assert document_payload["segments"][1]["textContent"] == "文氏图阴影区域表示 A∩B。"


def test_docx_parser_falls_back_to_vision_when_ocr_fails(tmp_path: Path):
    docx_path = tmp_path / "image-ocr-failed.docx"
    document = Document()
    document.add_heading("图表章节", level=1)
    document.add_picture(BytesIO(PNG_BYTES), width=Inches(1))
    document.save(docx_path)
    vision_client = FakeVisionClient([VisionResult(segment_type="image_caption", text="函数变化趋势图")])

    result = DocxParser(ocr_client=RaisingOcrClient(), vision_client=vision_client).parse(docx_path)

    assert result.status == "succeeded"
    document_payload = assert_valid_normalized_document(result.normalized_document)
    assert len(vision_client.calls) == 1
    assert document_payload["segments"][1]["textContent"] == "函数变化趋势图"
    assert issue_codes(result) == ["docx.ocr_failed"]


def test_docx_parser_reports_empty_visual_result(tmp_path: Path):
    docx_path = tmp_path / "image-empty-visual.docx"
    document = Document()
    document.add_heading("图表章节", level=1)
    document.add_picture(BytesIO(PNG_BYTES), width=Inches(1))
    document.save(docx_path)

    result = DocxParser(vision_client=FakeVisionClient([])).parse(docx_path)

    assert result.status == "succeeded"
    document_payload = assert_valid_normalized_document(result.normalized_document)
    assert [segment["segmentType"] for segment in document_payload["segments"]] == ["docx_block_text"]
    assert issue_codes(result) == ["docx.visual_empty"]


def test_docx_parser_reports_empty_ocr_without_vision(tmp_path: Path):
    docx_path = tmp_path / "image-empty-ocr.docx"
    document = Document()
    document.add_heading("图表章节", level=1)
    document.add_picture(BytesIO(PNG_BYTES), width=Inches(1))
    document.save(docx_path)
    ocr_client = FakeOcrClient([OcrAssetResult(asset_id="unused", text="B\nA\nANB", boxes=[])])

    result = DocxParser(ocr_client=ocr_client).parse(docx_path)

    assert result.status == "succeeded"
    document_payload = assert_valid_normalized_document(result.normalized_document)
    assert [segment["segmentType"] for segment in document_payload["segments"]] == ["docx_block_text"]
    assert issue_codes(result) == ["docx.ocr_empty"]


def test_srt_parser_extracts_captions_and_validates_schema(tmp_path: Path):
    srt_path = tmp_path / "captions.srt"
    srt_path.write_text(
        "\n".join(
            [
                "1",
                "00:00:01,000 --> 00:00:03,000",
                "First caption",
                "",
                "2",
                "00:00:04.250 --> 00:00:05.750",
                "Second",
                "caption",
                "",
            ]
        ),
        encoding="utf-8",
    )

    result = parse_resource("srt", srt_path)

    assert isinstance(result, ParserResult)
    assert result.status == "succeeded"
    document = assert_valid_normalized_document(result.normalized_document)
    assert document["resourceType"] == "srt"
    assert document["segments"] == [
        {
            "segmentKey": "srt-c1",
            "segmentType": "video_caption",
            "orderNo": 1,
            "textContent": "First caption",
            "startSec": 1,
            "endSec": 3,
        },
        {
            "segmentKey": "srt-c2",
            "segmentType": "video_caption",
            "orderNo": 2,
            "textContent": "Second\ncaption",
            "startSec": 4,
            "endSec": 6,
        },
    ]


def test_srt_parser_rejects_empty_or_invalid_timeline(tmp_path: Path):
    empty_srt_path = tmp_path / "empty.srt"
    empty_srt_path.write_text("", encoding="utf-8")
    empty_result = SrtParser().parse(empty_srt_path)
    assert empty_result.status == "failed"
    assert issue_codes(empty_result) == ["srt.caption_empty"]

    invalid_srt_path = tmp_path / "invalid.srt"
    invalid_srt_path.write_text(
        "1\n00:00:03,000 --> 00:00:01,000\nBad timeline\n",
        encoding="utf-8",
    )
    invalid_result = SrtParser().parse(invalid_srt_path)
    assert invalid_result.status == "failed"
    assert issue_codes(invalid_result) == ["srt.timeline_invalid"]


def test_mp4_parser_returns_asr_not_configured_issue(tmp_path: Path):
    video_path = tmp_path / "lecture.mp4"
    video_path.write_bytes(b"\x00\x00\x00\x18ftypmp42")

    result = VideoParser().parse(video_path)

    assert result.status == "failed"
    assert result.normalized_document is None
    assert issue_codes(result) == ["mp4.asr_not_configured"]


def test_first_edition_what_is_set_documents_parse_high_fidelity_under_75_seconds():
    asset_dir = ROOT / "local_assets/first-edition/what-is-set"
    required_files = [
        asset_dir / "knowlink-demo-handout.pdf",
        asset_dir / "knowlink-demo-slides.pptx",
        asset_dir / "knowlink-demo-notes.docx",
    ]
    if not all(path.exists() for path in required_files):
        pytest.skip("Local first-edition what-is-set assets are not available.")

    manifest = json.loads((ROOT / "server/seeds/demo_assets_manifest.json").read_text(encoding="utf-8"))
    for item in manifest["assets"]:
        path = asset_dir / item["relativePath"]
        assert path.exists(), path
        assert path.stat().st_size == item["sizeBytes"]
        assert f"sha256:{hashlib.sha256(path.read_bytes()).hexdigest()}" == item["checksum"]

    parse_inputs = [
        ("pdf", asset_dir / "knowlink-demo-handout.pdf"),
        ("pptx", asset_dir / "knowlink-demo-slides.pptx"),
        ("docx", asset_dir / "knowlink-demo-notes.docx"),
    ]
    started_at = time.perf_counter()
    results = {resource_type: parse_resource(resource_type, path) for resource_type, path in parse_inputs}
    elapsed_sec = time.perf_counter() - started_at

    assert elapsed_sec < 75
    for result in results.values():
        assert result.status == "succeeded"
        assert issue_codes(result) == []
        assert_valid_normalized_document(result.normalized_document)

    pdf_segments = results["pdf"].normalized_document["segments"]  # type: ignore[index]
    assert [segment["pageNo"] for segment in pdf_segments if segment["segmentType"] == "pdf_page_text"] == list(
        range(1, 11)
    )
    pdf_caption = next(
        segment for segment in pdf_segments if segment["segmentType"] == "image_caption" and segment["pageNo"] == 8
    )
    assert "文氏图" in pdf_caption["textContent"]
    assert "元素 a" in pdf_caption["textContent"]

    pptx_segments = results["pptx"].normalized_document["segments"]  # type: ignore[index]
    assert [segment["slideNo"] for segment in pptx_segments if segment["segmentType"] == "ppt_slide_text"] == list(
        range(1, 11)
    )
    pptx_caption = next(
        segment for segment in pptx_segments if segment["segmentType"] == "image_caption" and segment["slideNo"] == 8
    )
    assert "文氏图" in pptx_caption["textContent"]
    assert "元素 a" in pptx_caption["textContent"]

    docx_text = "\n".join(
        segment["textContent"] for segment in results["docx"].normalized_document["segments"]  # type: ignore[index]
    )
    assert "集合论基础测试题" in docx_text
    assert "单项选择题" in docx_text
    assert "参考答案" in docx_text


def test_real_parse_outputs_have_no_abnormal_chars_or_same_location_duplicates():
    output_dir = os.getenv("KNOWLINK_REAL_PARSE_OUTPUT_DIR")
    if not output_dir:
        pytest.skip("Set KNOWLINK_REAL_PARSE_OUTPUT_DIR to run local real fixture assertions.")

    paths = sorted(Path(output_dir).glob("*.json"))
    assert paths
    abnormal_re = re.compile(r"[\u0b80-\u0bff\ufffd\uffff\x00-\x08\x0b\x0c\x0e-\x1f\x7f-\x9f]")

    for path in paths:
        payload = json.loads(path.read_text(encoding="utf-8"))
        if not isinstance(payload, dict):
            continue
        segments = (payload.get("normalizedDocument") or {}).get("segments") or []
        seen_by_location: dict[tuple[str, object], list[str]] = {}
        for segment in segments:
            text = str(segment.get("textContent", ""))
            assert abnormal_re.search(text) is None, f"{path.name} contains abnormal chars: {text!r}"

            if "pageNo" in segment:
                location = ("pageNo", segment["pageNo"])
            elif "slideNo" in segment:
                location = ("slideNo", segment["slideNo"])
            else:
                continue

            normalized = compact_text(text)
            if len(normalized) < 4:
                continue
            existing = seen_by_location.setdefault(location, [])
            assert all(
                normalized != item and normalized not in item and item not in normalized for item in existing
            ), f"{path.name} has duplicate text at {location}: {text!r}"
            existing.append(normalized)


def test_vivo_vision_client_uses_batch_multimodal_chat_request(monkeypatch):
    captured = {}

    class FakeResponse:
        def __enter__(self):
            return self

        def __exit__(self, exc_type, exc, traceback):
            return None

        def read(self):
            return json.dumps(
                {
                    "choices": [
                        {
                            "message": {
                                "content": json.dumps(
                                    {
                                        "segments": [
                                            {
                                                "assetId": "pdf-p1",
                                                "segmentType": "image_caption",
                                                "textContent": "Venn 图展示 A∩B。",
                                            },
                                            {"assetId": "pdf-p2", "segmentType": "formula", "textContent": "A∩B"},
                                        ]
                                    },
                                    ensure_ascii=False,
                                )
                            }
                        }
                    ]
                },
                ensure_ascii=False,
            ).encode("utf-8")

    def fake_urlopen(request, timeout):
        captured["url"] = request.full_url
        captured["headers"] = dict(request.header_items())
        captured["body"] = request.data.decode("utf-8")
        captured["timeout"] = timeout
        return FakeResponse()

    monkeypatch.setattr("urllib.request.urlopen", fake_urlopen)

    client = VivoVisionClient(
        app_id="2026764332",
        app_key="fake-key",
        base_url="https://example.invalid",
        model="unused",
    )
    results = client.analyze_images(
        [
            VisualAsset(
                asset_id="pdf-p1",
                image_bytes=b"png-bytes-1",
                mime_type="image/png",
                location={"pageNo": 1},
                hint="pdf_page_visual",
            ),
            VisualAsset(
                asset_id="pdf-p2",
                image_bytes=b"png-bytes-2",
                mime_type="image/png",
                location={"pageNo": 2},
                hint="pdf_page_visual",
            ),
        ],
        resource_type="pdf",
        document_context="第 1 页：集合",
    )

    body = json.loads(captured["body"])
    assert captured["url"].startswith("https://example.invalid/v1/chat/completions?request_id=")
    assert captured["headers"]["Authorization"] == "Bearer fake-key"
    assert captured["headers"]["Content-type"] == "application/json; charset=utf-8"
    assert captured["timeout"] == 20.0
    assert body["model"] == "unused"
    prompt = body["messages"][0]["content"][0]["text"]
    assert "assetId" in prompt
    assert "只输出上下文缺失" in prompt
    assert "Markdown table" in prompt
    assert "公式尽量内联" in prompt
    assert body["messages"][0]["content"][2]["image_url"]["url"].startswith("data:image/png;base64,")
    assert body["messages"][0]["content"][4]["image_url"]["url"].startswith("data:image/png;base64,")
    assert results == [
        VisionAssetResult(asset_id="pdf-p1", segment_type="image_caption", text="Venn 图展示 A∩B。"),
        VisionAssetResult(asset_id="pdf-p2", segment_type="formula", text="A∩B"),
    ]


def test_vivo_vision_client_falls_back_to_image_caption_for_plain_response(monkeypatch):
    class FakeResponse:
        def __enter__(self):
            return self

        def __exit__(self, exc_type, exc, traceback):
            return None

        def read(self):
            return json.dumps({"choices": [{"message": {"content": "这是一张集合关系图。"}}]}, ensure_ascii=False).encode(
                "utf-8"
            )

    monkeypatch.setattr("urllib.request.urlopen", lambda request, timeout: FakeResponse())

    client = VivoVisionClient(
        app_id="2026764332",
        app_key="fake-key",
        base_url="https://example.invalid/v1",
        model="Doubao-Seed-2.0-mini",
    )
    results = client.analyze_image(
        b"png-bytes",
        mime_type="image/png",
        resource_type="pptx",
        location={"slideNo": 3},
        hint="pptx_shape_visual",
    )

    assert results == [VisionResult(segment_type="image_caption", text="这是一张集合关系图。")]


def test_vivo_vision_client_retries_multi_image_errors_as_single_requests(monkeypatch):
    calls = {"count": 0}

    class FakeResponse:
        def __enter__(self):
            return self

        def __exit__(self, exc_type, exc, traceback):
            return None

        def read(self):
            return json.dumps(
                {
                    "choices": [
                        {
                            "message": {
                                "content": json.dumps(
                                    {"segments": [{"segmentType": "ocr_text", "textContent": "single ok"}]},
                                    ensure_ascii=False,
                                )
                            }
                        }
                    ]
                },
                ensure_ascii=False,
            ).encode("utf-8")

    def fake_urlopen(request, timeout):
        calls["count"] += 1
        if calls["count"] == 1:
            raise OSError("too many image_url content items")
        return FakeResponse()

    monkeypatch.setattr("urllib.request.urlopen", fake_urlopen)
    client = VivoVisionClient(
        app_id="2026764332",
        app_key="fake-key",
        base_url="https://example.invalid/v1",
        model="Doubao-Seed-2.0-mini",
    )

    results = client.analyze_images(
        [
            VisualAsset(asset_id="a1", image_bytes=b"one", mime_type="image/png", location={"pageNo": 1}),
            VisualAsset(asset_id="a2", image_bytes=b"two", mime_type="image/png", location={"pageNo": 2}),
        ],
        resource_type="pdf",
    )

    assert calls["count"] == 3
    assert results == [
        VisionAssetResult(asset_id="a1", segment_type="ocr_text", text="single ok"),
        VisionAssetResult(asset_id="a2", segment_type="ocr_text", text="single ok"),
    ]


def test_vivo_ocr_client_uses_general_recognition_with_pos_2(monkeypatch):
    captured = {}

    class FakeResponse:
        def __enter__(self):
            return self

        def __exit__(self, exc_type, exc, traceback):
            return None

        def read(self):
            return json.dumps(
                {
                    "error_code": 0,
                    "error_msg": "succ",
                    "result": {
                        "OCR": [
                            {
                                "words": "第一行",
                                "location": {
                                    "top_left": {"x": 0.1, "y": 0.2},
                                    "top_right": {"x": 0.4, "y": 0.2},
                                    "down_left": {"x": 0.1, "y": 0.3},
                                    "down_right": {"x": 0.4, "y": 0.3},
                                },
                            },
                            {
                                "words": "第二行",
                                "location": {
                                    "top_left": {"x": 0.1, "y": 0.4},
                                    "top_right": {"x": 0.4, "y": 0.4},
                                    "down_left": {"x": 0.1, "y": 0.5},
                                    "down_right": {"x": 0.4, "y": 0.5},
                                },
                            },
                        ],
                        "angle": 0,
                    },
                },
                ensure_ascii=False,
            ).encode("utf-8")

    def fake_urlopen(request, timeout):
        captured["url"] = request.full_url
        captured["headers"] = dict(request.header_items())
        captured["body"] = request.data.decode("utf-8")
        captured["timeout"] = timeout
        return FakeResponse()

    monkeypatch.setattr("urllib.request.urlopen", fake_urlopen)
    client = VivoOcrClient(
        app_key="fake-key",
        business_id="aigc2026764332",
        base_url="https://example.invalid",
        timeout_sec=7,
    )

    results = client.recognize_images(
        [
            type(
                "Asset",
                (),
                {
                    "asset_id": "pdf-p1",
                    "image_bytes": b"png-bytes",
                    "mime_type": "image/png",
                    "location": {"pageNo": 1},
                    "hint": "pdf_page_visual",
                },
            )()
        ],
        resource_type="pdf",
    )

    body = dict(item.split("=", 1) for item in captured["body"].split("&"))
    assert captured["url"].startswith("https://example.invalid/ocr/general_recognition?requestId=")
    assert captured["headers"]["Authorization"] == "Bearer fake-key"
    assert captured["headers"]["Content-type"] == "application/x-www-form-urlencoded"
    assert captured["timeout"] == 7
    assert body["pos"] == "2"
    assert body["businessid"] == "aigc2026764332"
    assert results == [
        OcrAssetResult(
            asset_id="pdf-p1",
            text="第一行\n第二行",
            boxes=[
                OcrBox(text="第一行", x=0.1, y=0.2, w=0.30000000000000004, h=0.09999999999999998),
                OcrBox(text="第二行", x=0.1, y=0.4, w=0.30000000000000004, h=0.09999999999999998),
            ],
        )
    ]
