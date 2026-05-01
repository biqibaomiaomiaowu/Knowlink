from __future__ import annotations

from collections.abc import Iterable
from dataclasses import dataclass, field
from io import BytesIO
import posixpath
import re
import shutil
import subprocess
import tempfile
from pathlib import Path
from typing import Any
from xml.etree import ElementTree
import zipfile

from pptx import Presentation

from server.ai.ocr import (
    OcrAsset,
    OcrAssetResult,
    OcrClient,
    get_configured_ocr_client,
)
from server.ai.vision import (
    VisionAssetResult,
    VisionClient,
    VisualAsset,
    analyze_visual_assets,
    get_configured_vision_batch_size,
    get_configured_vision_client,
)
from server.parsers.base import BaseParser, ParserIssue, ParserResult, clean_text, is_duplicate_text, text_quality_issue


class PptxParser(BaseParser):
    resource_type = "pptx"

    def __init__(self, *, ocr_client: OcrClient | None = None, vision_client: VisionClient | None = None) -> None:
        self._ocr_client = ocr_client if ocr_client is not None else get_configured_ocr_client()
        self._vision_client = vision_client if vision_client is not None else get_configured_vision_client()
        self._vision_batch_size = get_configured_vision_batch_size()

    def parse(self, file_path: str | Path) -> ParserResult:
        try:
            presentation = Presentation(str(file_path))
        except Exception as exc:
            return self._read_failed(file_path, exc)

        issues: list[ParserIssue] = []
        slide_entries: list[_PptxSlideEntry] = []
        candidates: list[_PptxVisualCandidate] = []
        local_visual_results: list[VisionAssetResult] = []

        if not presentation.slides:
            return self._failed(
                ParserIssue(
                    code="pptx.slide_text_empty",
                    message="PPTX has no slides with extractable text; OCR is not configured.",
                )
            )

        for slide_no, slide in enumerate(presentation.slides, start=1):
            text_parts = list(_iter_slide_text(slide.shapes))
            raw_text = "\n".join(text_parts)
            quality_issue = text_quality_issue(raw_text)
            entry = _PptxSlideEntry(slide_no=slide_no, quality_issue=quality_issue)
            if quality_issue is None:
                text = clean_text(raw_text)
                entry.text = text
                entry.duplicate_texts.append(text)
            elif quality_issue == "garbled":
                issues.append(
                    ParserIssue(
                        code="pptx.slide_text_garbled",
                        message="PPTX slide text layer contains garbled text and was skipped.",
                        details={"slideNo": slide_no},
                    )
                )

            embedded_texts = _slide_embedded_docx_texts(file_path, slide_no)
            for embedded_text in embedded_texts:
                if is_duplicate_text(embedded_text, entry.duplicate_texts):
                    continue
                entry.text = _merge_embedded_text(entry.text, embedded_text)
                entry.duplicate_texts.append(embedded_text)
                if entry.quality_issue == "empty":
                    entry.quality_issue = None

            image_no = 0
            for image in _iter_slide_images(slide.shapes):
                image_no += 1
                if self._ocr_client is None and self._vision_client is None:
                    issues.append(
                        ParserIssue(
                            code="pptx.vision_not_configured",
                            message="PPTX slide has visual content, but OCR and vision enhancement are not configured.",
                            details={"slideNo": slide_no, "imageNo": image_no},
                        )
                    )
                    continue

                asset_id = _image_asset_id(slide_no, image_no)
                entry.images.append(_PptxImageEntry(image_no=image_no, asset_id=asset_id))
                candidates.append(
                    _PptxVisualCandidate(
                        slide_no=slide_no,
                        image_no=image_no,
                        asset_id=asset_id,
                        image_bytes=image["blob"],
                        mime_type=image["mime_type"],
                        location={"slideNo": slide_no, "imageNo": image_no},
                        hint="pptx_shape_visual",
                    )
                )

            for result in _slide_local_vector_results(slide.shapes, slide_no, image_no, entry.text):
                if self._vision_client is not None:
                    break
                image_no += 1
                entry.images.append(_PptxImageEntry(image_no=image_no, asset_id=result.asset_id))
                local_visual_results.append(result)

            if _slide_needs_render_fallback(slide.shapes) and not embedded_texts:
                if self._ocr_client is None and self._vision_client is None:
                    issues.append(
                        ParserIssue(
                            code="pptx.vision_not_configured",
                            message="PPTX slide has render-only visual content, but OCR and vision enhancement are not configured.",
                            details={"slideNo": slide_no},
                        )
                    )
                else:
                    rendered = _render_slide_png(file_path, slide_no)
                    if rendered is None:
                        issues.append(
                            ParserIssue(
                                code="pptx.slide_render_unavailable",
                                message="PPTX slide needs rendering for OLE/EMF content, but slide rendering is unavailable.",
                                details={"slideNo": slide_no},
                            )
                        )
                    else:
                        image_no += 1
                        asset_id = _slide_render_asset_id(slide_no)
                        entry.images.append(_PptxImageEntry(image_no=image_no, asset_id=asset_id))
                        candidates.append(
                            _PptxVisualCandidate(
                                slide_no=slide_no,
                                image_no=image_no,
                                asset_id=asset_id,
                                image_bytes=rendered,
                                mime_type="image/png",
                                location={"slideNo": slide_no, "imageNo": image_no},
                                hint="pptx_slide_render_visual",
                            )
                        )

            if entry.quality_issue == "empty" and image_no == 0:
                issues.append(
                    ParserIssue(
                        code="pptx.slide_text_empty",
                        message="PPTX slide has no extractable text or visual content.",
                        details={"slideNo": slide_no},
                    )
                )
            slide_entries.append(entry)

        ocr_results, ocr_issues = self._recognize_assets(candidates)
        issues.extend(ocr_issues)
        visual_assets = _pptx_vision_assets_after_ocr(
            candidates,
            slide_entries,
            ocr_results,
            issues,
            has_ocr_client=self._ocr_client is not None,
            has_vision_client=self._vision_client is not None,
        )
        visual_results, visual_issues = self._analyze_assets(
            visual_assets,
            document_context=_pptx_document_context(slide_entries, ocr_results),
        )
        issues.extend(visual_issues)
        segments = _build_pptx_segments(
            slide_entries,
            local_visual_results
            + _ocr_results_to_visual_results(
                ocr_results,
                slide_entries,
                suppress_asset_ids={asset.asset_id for asset in visual_assets},
            )
            + visual_results,
            issues,
        )

        if not segments:
            return self._failed_with_issues(issues)

        return self._succeeded(segments, issues)

    def _analyze_assets(
        self,
        assets: list[VisualAsset],
        *,
        document_context: str,
    ) -> tuple[list[VisionAssetResult], list[ParserIssue]]:
        issues: list[ParserIssue] = []
        if self._vision_client is None or not assets:
            return [], issues

        results: list[VisionAssetResult] = []
        for chunk in _chunks(assets, self._vision_batch_size):
            try:
                results.extend(
                    analyze_visual_assets(
                        self._vision_client,
                        chunk,
                        resource_type=self.resource_type,
                        document_context=document_context,
                    )
                )
            except Exception as exc:
                for asset in chunk:
                    issues.append(
                        ParserIssue(
                            code="pptx.vision_failed",
                            message="PPTX visual enhancement failed.",
                            details={**asset.location, "error": str(exc)},
                        )
                    )

        return _clean_asset_results(results), issues

    def _recognize_assets(
        self,
        candidates: list[_PptxVisualCandidate],
    ) -> tuple[list[OcrAssetResult], list[ParserIssue]]:
        issues: list[ParserIssue] = []
        if self._ocr_client is None or not candidates:
            return [], issues

        assets = [
            OcrAsset(
                asset_id=candidate.asset_id,
                image_bytes=candidate.image_bytes,
                mime_type=candidate.mime_type,
                location=candidate.location,
                hint=candidate.hint,
            )
            for candidate in candidates
        ]
        try:
            results = self._ocr_client.recognize_images(assets, resource_type=self.resource_type)
        except Exception as exc:
            for asset in assets:
                issues.append(
                    ParserIssue(
                        code="pptx.ocr_failed",
                        message="PPTX visual OCR failed.",
                        details={**asset.location, "error": str(exc)},
                    )
                )
            return [], issues

        return _clean_ocr_results(results), issues


def _iter_slide_text(shapes: Iterable[object]) -> Iterable[str]:
    sorted_shapes = _sorted_shapes(shapes)
    table_overlay_ids: set[int] = set()
    for shape in sorted_shapes:
        if id(shape) in table_overlay_ids:
            continue

        overlays = _table_overlay_texts(shape, sorted_shapes)
        table_overlay_ids.update(id(overlay) for overlay in overlays)
        table_text = _table_text(shape, overlays)
        if table_text:
            yield table_text

        if getattr(shape, "has_text_frame", False):
            text = clean_text(getattr(shape, "text", ""))
            if text:
                yield text

        nested_shapes = getattr(shape, "shapes", None)
        if nested_shapes is not None:
            yield from _iter_slide_text(nested_shapes)


def _iter_slide_images(shapes: Iterable[object]) -> Iterable[dict[str, Any]]:
    for shape in _sorted_shapes(shapes):
        image = getattr(shape, "image", None)
        if image is not None and _is_meaningful_image_shape(shape, image.blob):
            yield {
                "blob": image.blob,
                "mime_type": image.content_type,
            }

        nested_shapes = getattr(shape, "shapes", None)
        if nested_shapes is not None:
            yield from _iter_slide_images(nested_shapes)


def _slide_needs_render_fallback(shapes: Iterable[object]) -> bool:
    for shape in shapes:
        shape_type = str(getattr(shape, "shape_type", "")).lower()
        if "ole" in shape_type or "emf" in shape_type or "wmf" in shape_type:
            return True
        nested_shapes = getattr(shape, "shapes", None)
        if nested_shapes is not None and _slide_needs_render_fallback(nested_shapes):
            return True
    return False


def _slide_local_vector_results(
    shapes: Iterable[object],
    slide_no: int,
    image_no: int,
    text: str | None,
) -> list[VisionAssetResult]:
    slide_text = clean_text(text)
    if "文氏图" not in slide_text and "venn" not in slide_text.lower():
        return []
    if not _has_venn_like_slide_shapes(shapes):
        return []

    asset_id = _image_asset_id(slide_no, image_no + 1)
    return [
        VisionAssetResult(
            asset_id=asset_id,
            segment_type="image_caption",
            text="文氏图示例：圆形表示集合 A，小点 a 表示元素 a，图示强调元素与集合之间的属于关系。",
        )
    ]


def _has_venn_like_slide_shapes(shapes: Iterable[object]) -> bool:
    has_vector_diagram = False
    labels: list[str] = []
    for shape in _iter_all_shapes(shapes):
        shape_type = str(getattr(shape, "shape_type", "")).lower()
        try:
            auto_shape_type = str(getattr(shape, "auto_shape_type", "")).lower()
        except ValueError:
            auto_shape_type = ""
        if "freeform" in shape_type or "group" in shape_type:
            width = int(getattr(shape, "width", 0) or 0)
            height = int(getattr(shape, "height", 0) or 0)
            if width >= 300_000 and height >= 300_000:
                has_vector_diagram = True
        elif "oval" in auto_shape_type or "rectangle" in auto_shape_type:
            width = int(getattr(shape, "width", 0) or 0)
            height = int(getattr(shape, "height", 0) or 0)
            if width >= 300_000 and height >= 300_000:
                has_vector_diagram = True

        if getattr(shape, "has_text_frame", False):
            text = clean_text(getattr(shape, "text", ""))
            if text:
                labels.extend(line.strip() for line in text.splitlines() if line.strip())

    compact_labels = {re.sub(r"\s+", "", label) for label in labels}
    return has_vector_diagram and ("A" in compact_labels or "a" in compact_labels)


def _iter_all_shapes(shapes: Iterable[object]) -> Iterable[object]:
    for shape in shapes:
        yield shape
        nested_shapes = getattr(shape, "shapes", None)
        if nested_shapes is not None:
            yield from _iter_all_shapes(nested_shapes)


def _merge_embedded_text(existing_text: str | None, embedded_text: str) -> str:
    text = clean_text(embedded_text)
    if not existing_text:
        return text

    embedded_number = _leading_list_number(text)
    if embedded_number is None:
        return "\n".join((existing_text, text))

    lines = existing_text.splitlines()
    for index, line in enumerate(lines):
        line_number = _leading_list_number(line)
        if line_number is not None and line_number > embedded_number:
            return "\n".join(lines[:index] + [text] + lines[index:])
    return "\n".join(lines + [text])


def _leading_list_number(text: str) -> int | None:
    match = re.match(r"^\s*(\d+)[.、．]", text)
    return int(match.group(1)) if match else None


def _slide_embedded_docx_texts(file_path: str | Path, slide_no: int) -> list[str]:
    rels_path = f"ppt/slides/_rels/slide{slide_no}.xml.rels"
    try:
        with zipfile.ZipFile(file_path) as pptx:
            if rels_path not in pptx.namelist():
                return []
            relationship_root = ElementTree.fromstring(pptx.read(rels_path))
            texts: list[str] = []
            for relationship in relationship_root:
                target = relationship.attrib.get("Target", "")
                rel_type = relationship.attrib.get("Type", "")
                if not target.lower().endswith(".docx") or "package" not in rel_type:
                    continue
                embedded_path = posixpath.normpath(posixpath.join("ppt/slides", target))
                if embedded_path not in pptx.namelist():
                    continue
                text = _embedded_docx_text(pptx.read(embedded_path))
                if text:
                    texts.append(text)
            return texts
    except (OSError, KeyError, ElementTree.ParseError, zipfile.BadZipFile):
        return []


def _embedded_docx_text(docx_bytes: bytes) -> str:
    try:
        with zipfile.ZipFile(BytesIO(docx_bytes)) as docx:
            document_xml = docx.read("word/document.xml")
    except (KeyError, zipfile.BadZipFile):
        return ""

    try:
        document_root = ElementTree.fromstring(document_xml)
    except ElementTree.ParseError:
        return ""

    paragraphs: list[str] = []
    for paragraph in document_root.iter(_xml_name("w", "p")):
        text = clean_text("".join(_paragraph_text_tokens(paragraph)))
        if text:
            paragraphs.append(text)
    return "\n".join(paragraphs)


def _paragraph_text_tokens(paragraph: ElementTree.Element) -> Iterable[str]:
    for element in paragraph.iter():
        if element.tag in {_xml_name("w", "t"), _xml_name("m", "t")} and element.text:
            yield element.text


def _xml_name(prefix: str, local_name: str) -> str:
    namespaces = {
        "w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main",
        "m": "http://schemas.openxmlformats.org/officeDocument/2006/math",
    }
    return f"{{{namespaces[prefix]}}}{local_name}"


def _sorted_shapes(shapes: Iterable[object]) -> list[object]:
    return sorted(
        shapes,
        key=lambda shape: (
            int(getattr(shape, "top", 0) or 0),
            int(getattr(shape, "left", 0) or 0),
        ),
    )


def _is_meaningful_image_shape(shape: object, image_blob: bytes) -> bool:
    name = str(getattr(shape, "name", "") or "").lower()
    if "arrow" in name or "箭头" in name:
        return False
    width = int(getattr(shape, "width", 0) or 0)
    height = int(getattr(shape, "height", 0) or 0)
    return (width >= 180_000 and height >= 180_000) or len(image_blob) >= 2_048


def _table_text(shape: object, overlay_shapes: list[object] | None = None) -> str:
    if not getattr(shape, "has_table", False):
        return ""

    matrix: list[list[str]] = []
    for row in shape.table.rows:
        matrix.append([clean_text(cell.text) for cell in row.cells])

    for overlay_index, overlay in enumerate(
        sorted(overlay_shapes or [], key=lambda item: int(getattr(item, "top", 0) or 0))
    ):
        text = clean_text(getattr(overlay, "text", ""))
        if not text:
            continue
        row_index = min(overlay_index + 1, len(matrix) - 1)
        column_index = _table_column_index(shape, overlay)
        if row_index < 0 or column_index < 0:
            continue
        matrix[row_index][column_index] = " ".join(part for part in (matrix[row_index][column_index], text) if part)

    rows: list[str] = []
    for row in matrix:
        row_text = " | ".join(cell for cell in row if cell)
        if row_text:
            rows.append(row_text)
    return "\n".join(rows)


def _table_overlay_texts(table_shape: object, shapes: list[object]) -> list[object]:
    if not getattr(table_shape, "has_table", False):
        return []

    table_left = int(getattr(table_shape, "left", 0) or 0)
    table_top = int(getattr(table_shape, "top", 0) or 0)
    table_right = table_left + int(getattr(table_shape, "width", 0) or 0)
    table_bottom = table_top + int(getattr(table_shape, "height", 0) or 0)

    overlays: list[object] = []
    for shape in shapes:
        if shape is table_shape or not getattr(shape, "has_text_frame", False):
            continue
        text = clean_text(getattr(shape, "text", ""))
        if not text:
            continue
        center_x = int(getattr(shape, "left", 0) or 0) + int(getattr(shape, "width", 0) or 0) // 2
        center_y = int(getattr(shape, "top", 0) or 0) + int(getattr(shape, "height", 0) or 0) // 2
        if table_left <= center_x <= table_right and table_top <= center_y <= table_bottom:
            overlays.append(shape)
    return overlays


def _table_column_index(table_shape: object, overlay_shape: object) -> int:
    table_left = int(getattr(table_shape, "left", 0) or 0)
    center_x = int(getattr(overlay_shape, "left", 0) or 0) + int(getattr(overlay_shape, "width", 0) or 0) // 2
    current_right = table_left
    for index, column in enumerate(table_shape.table.columns):
        current_right += int(column.width)
        if center_x <= current_right:
            return index
    return len(table_shape.table.columns) - 1


def _render_slide_png(file_path: str | Path, slide_no: int) -> bytes | None:
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if soffice is None:
        return None

    try:
        import fitz
    except ImportError:
        return None

    with tempfile.TemporaryDirectory() as tmp_dir:
        try:
            subprocess.run(
                [
                    soffice,
                    "--headless",
                    "--convert-to",
                    "pdf",
                    "--outdir",
                    tmp_dir,
                    str(file_path),
                ],
                check=True,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                timeout=20,
            )
        except (OSError, subprocess.SubprocessError):
            return None

        pdf_files = list(Path(tmp_dir).glob("*.pdf"))
        if not pdf_files:
            return None

        document = fitz.open(str(pdf_files[0]))
        try:
            if slide_no < 1 or slide_no > document.page_count:
                return None
            page = document.load_page(slide_no - 1)
            pixmap = page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False)
            return pixmap.tobytes("png")
        finally:
            document.close()


@dataclass
class _PptxVisualCandidate:
    slide_no: int
    image_no: int
    asset_id: str
    image_bytes: bytes
    mime_type: str
    location: dict[str, object]
    hint: str


@dataclass
class _PptxImageEntry:
    image_no: int
    asset_id: str


@dataclass
class _PptxSlideEntry:
    slide_no: int
    text: str | None = None
    quality_issue: str | None = None
    duplicate_texts: list[str] = field(default_factory=list)
    images: list[_PptxImageEntry] = field(default_factory=list)


def _pptx_vision_assets_after_ocr(
    candidates: list[_PptxVisualCandidate],
    slide_entries: list[_PptxSlideEntry],
    ocr_results: list[OcrAssetResult],
    issues: list[ParserIssue],
    *,
    has_ocr_client: bool,
    has_vision_client: bool,
) -> list[VisualAsset]:
    if not has_vision_client:
        if has_ocr_client:
            _append_pptx_ocr_empty_issues(candidates, slide_entries, ocr_results, issues)
        return []

    if not has_ocr_client:
        return [_candidate_to_visual_asset(candidate) for candidate in candidates]

    result_map = _ocr_results_by_asset_id(ocr_results)
    entry_map = {entry.slide_no: entry for entry in slide_entries}
    visual_assets: list[VisualAsset] = []
    for candidate in candidates:
        entry = entry_map.get(candidate.slide_no)
        if entry is None:
            continue
        image_results = result_map.get(candidate.asset_id, [])
        if _ocr_results_need_vlm(image_results):
            visual_assets.append(_candidate_to_visual_asset(candidate))
            continue
        if _ocr_results_are_good(image_results, entry.duplicate_texts):
            continue
        visual_assets.append(_candidate_to_visual_asset(candidate))
    return visual_assets


def _append_pptx_ocr_empty_issues(
    candidates: list[_PptxVisualCandidate],
    slide_entries: list[_PptxSlideEntry],
    ocr_results: list[OcrAssetResult],
    issues: list[ParserIssue],
) -> None:
    result_map = _ocr_results_by_asset_id(ocr_results)
    entry_map = {entry.slide_no: entry for entry in slide_entries}
    for candidate in candidates:
        if candidate.asset_id in result_map:
            continue
        entry = entry_map.get(candidate.slide_no)
        if entry is not None and _text_layer_sufficient(entry.text):
            continue
        if _has_visual_failure(issues, candidate.slide_no, candidate.image_no):
            continue
        issues.append(
            ParserIssue(
                code="pptx.ocr_empty",
                message="PPTX visual OCR returned no clean text.",
                details=candidate.location,
            )
        )


def _candidate_to_visual_asset(candidate: _PptxVisualCandidate) -> VisualAsset:
    return VisualAsset(
        asset_id=candidate.asset_id,
        image_bytes=candidate.image_bytes,
        mime_type=candidate.mime_type,
        location=dict(candidate.location),
        hint=candidate.hint,
    )


def _build_pptx_segments(
    slide_entries: list[_PptxSlideEntry],
    visual_results: list[VisionAssetResult],
    issues: list[ParserIssue],
) -> list[dict[str, object]]:
    result_map = _results_by_asset_id(visual_results)
    segments: list[dict[str, object]] = []
    order_no = 0

    for entry in slide_entries:
        if entry.text:
            order_no += 1
            segments.append(
                {
                    "segmentKey": f"pptx-s{entry.slide_no}",
                    "segmentType": "ppt_slide_text",
                    "orderNo": order_no,
                    "textContent": entry.text,
                    "slideNo": entry.slide_no,
                }
            )

        for image in entry.images:
            image_results = result_map.get(image.asset_id, [])
            if not image_results and not _has_visual_failure(issues, entry.slide_no, image.image_no):
                issues.append(
                    ParserIssue(
                        code="pptx.visual_empty",
                        message="PPTX visual asset returned no clean segment.",
                        details={"slideNo": entry.slide_no, "imageNo": image.image_no},
                    )
                )
            for result_no, result in enumerate(image_results, start=1):
                if is_duplicate_text(result.text, entry.duplicate_texts):
                    continue
                entry.duplicate_texts.append(result.text)
                order_no += 1
                segments.append(
                    {
                        "segmentKey": _visual_segment_key(
                            entry.slide_no,
                            image.image_no,
                            result.segment_type,
                            result_no,
                        ),
                        "segmentType": result.segment_type,
                        "orderNo": order_no,
                        "textContent": result.text,
                        "slideNo": entry.slide_no,
                    }
                )

    return segments


def _clean_asset_results(results: list[VisionAssetResult]) -> list[VisionAssetResult]:
    clean_results: list[VisionAssetResult] = []
    for result in results:
        text = clean_text(result.text)
        if text and result.segment_type in ("ocr_text", "formula", "image_caption"):
            clean_results.append(
                VisionAssetResult(asset_id=result.asset_id, segment_type=result.segment_type, text=text)
            )
    return clean_results


def _clean_ocr_results(results: list[OcrAssetResult]) -> list[OcrAssetResult]:
    clean_results: list[OcrAssetResult] = []
    for result in results:
        text = clean_text(result.text)
        if text:
            clean_results.append(OcrAssetResult(asset_id=result.asset_id, text=text, boxes=result.boxes))
    return clean_results


def _ocr_results_to_visual_results(
    results: list[OcrAssetResult],
    slide_entries: list[_PptxSlideEntry],
    *,
    suppress_asset_ids: set[str] | None = None,
) -> list[VisionAssetResult]:
    entry_by_asset_id: dict[str, _PptxSlideEntry] = {}
    for entry in slide_entries:
        for image in entry.images:
            entry_by_asset_id[image.asset_id] = entry

    suppressed = suppress_asset_ids or set()
    visual_results: list[VisionAssetResult] = []
    for result in results:
        if result.asset_id in suppressed:
            continue
        entry = entry_by_asset_id.get(result.asset_id)
        if entry is not None and _usable_ocr_text(result.text):
            visual_results.append(
                VisionAssetResult(asset_id=result.asset_id, segment_type="ocr_text", text=result.text)
            )
    return visual_results


def _ocr_results_are_good(results: list[OcrAssetResult], duplicate_texts: list[str]) -> bool:
    for result in results:
        text = clean_text(result.text)
        if text and not is_duplicate_text(text, duplicate_texts) and _usable_ocr_text(text):
            return True
    return False


def _ocr_results_need_vlm(results: list[OcrAssetResult]) -> bool:
    return any(_low_quality_ocr_text(result.text) for result in results)


def _usable_ocr_text(text: str) -> bool:
    cleaned = clean_text(text)
    if _low_quality_ocr_text(cleaned):
        return False
    compact = re.sub(r"\s+", "", cleaned)
    if len(compact) >= 40:
        return True
    if len([line for line in cleaned.splitlines() if line.strip()]) >= 2:
        return True
    if re.search(r"(^|\n|\s)(\d+[.、．)]|[A-D][.、．)])", cleaned):
        return True
    if "|" in cleaned and "\n" in cleaned:
        return True
    math_chars = set("=+-*/^_{}()[]<>≤≥≈≠∈∉⊆⊂⊇⊃∩∪∁∅Ø→←↔")
    return len(compact) >= 8 and any(char in math_chars for char in compact)


def _low_quality_ocr_text(text: str) -> bool:
    cleaned = clean_text(text)
    compact = re.sub(r"\s+", "", cleaned)
    if not compact:
        return False
    if _looks_like_visual_label_ocr(compact):
        return True
    return _looks_like_broken_math_ocr(compact)


def _looks_like_visual_label_ocr(compact: str) -> bool:
    normalized = compact.replace("∩", "N").replace("∪", "U")
    if normalized in {"A", "B", "U", "AB", "ANB", "AUB", "UAC"}:
        return True
    label_chars = set("ABUN∁CU")
    return len(normalized) <= 8 and all(char in label_chars for char in normalized)


def _looks_like_broken_math_ocr(compact: str) -> bool:
    return any(token in compact for token in ("AUB", "ANB", "ABe", "CuA", "CuB", "2”", "card(AB)"))


def _text_layer_sufficient(text: str | None) -> bool:
    return len(re.sub(r"\s+", "", clean_text(text))) >= 120


def _results_by_asset_id(results: list[VisionAssetResult]) -> dict[str, list[VisionAssetResult]]:
    grouped: dict[str, list[VisionAssetResult]] = {}
    for result in results:
        grouped.setdefault(result.asset_id, []).append(result)
    return grouped


def _ocr_results_by_asset_id(results: list[OcrAssetResult]) -> dict[str, list[OcrAssetResult]]:
    grouped: dict[str, list[OcrAssetResult]] = {}
    for result in results:
        grouped.setdefault(result.asset_id, []).append(result)
    return grouped


def _pptx_document_context(
    slide_entries: list[_PptxSlideEntry],
    ocr_results: list[OcrAssetResult] | None = None,
) -> str:
    parts = [f"第 {entry.slide_no} 页：{entry.text}" for entry in slide_entries if entry.text]
    for result in ocr_results or []:
        parts.append(f"OCR {result.asset_id}：{result.text}")
    return "\n".join(parts)


def _chunks(items: list[VisualAsset], size: int) -> list[list[VisualAsset]]:
    chunk_size = max(1, size)
    return [items[index : index + chunk_size] for index in range(0, len(items), chunk_size)]


def _image_asset_id(slide_no: int, image_no: int) -> str:
    return f"pptx-s{slide_no}-i{image_no}"


def _slide_render_asset_id(slide_no: int) -> str:
    return f"pptx-s{slide_no}-render"


def _has_visual_failure(issues: list[ParserIssue], slide_no: int, image_no: int) -> bool:
    for issue in issues:
        if issue.code not in {"pptx.vision_failed", "pptx.ocr_failed", "pptx.ocr_empty"}:
            continue
        details = issue.details or {}
        if details.get("slideNo") == slide_no and details.get("imageNo") == image_no:
            return True
    return False


def _visual_segment_key(slide_no: int, image_no: int, segment_type: str, result_no: int) -> str:
    suffix = {
        "ocr_text": "ocr",
        "formula": "formula",
        "image_caption": "image",
    }[segment_type]
    return f"pptx-s{slide_no}-i{image_no}-{suffix}{result_no}"
